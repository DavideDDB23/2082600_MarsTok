"""
Mars Operations — Styled PPTX builder v3
Dark space theme matching the Marp deck.
Run:  python3 build_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ───────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0f, 0x17, 0x2a)
BG_CARD = RGBColor(0x1e, 0x29, 0x3b)
BG_DARK = RGBColor(0x0b, 0x11, 0x1e)
BORDER  = RGBColor(0x33, 0x41, 0x55)
ORANGE  = RGBColor(0xf9, 0x73, 0x16)
ORNGL   = RGBColor(0xfb, 0x92, 0x3c)
WHITE   = RGBColor(0xff, 0xff, 0xff)
TEXT    = RGBColor(0xe2, 0xe8, 0xf0)
MUTED   = RGBColor(0x94, 0xa3, 0xb8)
RED     = RGBColor(0xef, 0x44, 0x44)
GREEN   = RGBColor(0x22, 0xc5, 0x5e)
AMBER   = RGBColor(0xf5, 0x9e, 0x0b)
CYAN    = RGBColor(0x06, 0xb6, 0xd4)
BLUE    = RGBColor(0x60, 0xa5, 0xfa)
PURPLE  = RGBColor(0xa8, 0x55, 0xf7)
CODE    = RGBColor(0x34, 0xd3, 0x99)

W, H = Inches(13.33), Inches(7.5)
FTR   = Inches(0.30)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

def hex6(c):
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"

def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def rect(slide, x, y, w, h, fill, line=None, lw=Pt(0.75)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s

def txt(slide, x, y, w, h, text, sz=Pt(15), bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, font="Calibri", italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size    = sz
    r.font.bold    = bold
    r.font.italic  = italic
    r.font.color.rgb = color
    r.font.name    = font
    return tb

def card(slide, x, y, w, h, accent=None):
    rect(slide, x, y, w, h, BG_CARD, BORDER)
    if accent:
        rect(slide, x, y, Inches(0.07), h, accent)

def footer(slide, pg):
    rect(slide, 0, H - FTR, W, FTR, BG_DARK)
    rect(slide, 0, H - FTR, W, Inches(0.018), BORDER)
    txt(slide, Inches(0.5), H - FTR + Inches(0.07), Inches(6), FTR - Inches(0.1),
        "Mars Operations  ·  Hackathon March 2026  ·  Student 2082600",
        sz=Pt(9), color=MUTED)
    txt(slide, W - Inches(0.88), H - FTR + Inches(0.07), Inches(0.68), FTR - Inches(0.1),
        str(pg), sz=Pt(10), bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)

def header(slide, title):
    rect(slide, 0, 0, W, Inches(0.09), ORANGE)
    rect(slide, 0, Inches(0.09), W, Inches(0.72), BG_CARD)
    txt(slide, Inches(0.5), Inches(0.11), Inches(12.4), Inches(0.58),
        title, sz=Pt(26), bold=True, color=ORANGE)
    rect(slide, Inches(0.5), Inches(0.84), Inches(12.33), Inches(0.025), ORANGE)

CTOP = Inches(0.97)
CBOT = H - FTR - Inches(0.08)

def add_table(slide, x, y, w, rows, col_ws=None, rh=Inches(0.40)):
    nr, nc = len(rows), len(rows[0])
    tbl = slide.shapes.add_table(nr, nc, x, y, w, rh * nr).table
    if col_ws:
        for i, cw in enumerate(col_ws): tbl.columns[i].width = cw
    else:
        cw = int(w / nc)
        for i in range(nc): tbl.columns[i].width = cw
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = rh
        for ci, ct in enumerate(row):
            cell = tbl.cell(ri, ci)
            tc   = cell._tc
            tcPr = tc.get_or_add_tcPr()
            bv = hex6(BG_CARD) if ri % 2 == 0 else "131d2e"
            sf = etree.SubElement(tcPr, qn('a:solidFill'))
            etree.SubElement(sf, qn('a:srgbClr')).set('val', bv)
            lb = etree.SubElement(tcPr, qn('a:lnB'), w="9525", cap="flat", cmpd="sng")
            sf2 = etree.SubElement(lb, qn('a:solidFill'))
            etree.SubElement(sf2, qn('a:srgbClr')).set('val', hex6(BORDER))
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p._p.clear()
            r = p.add_run()
            r.text = ct
            r.font.name = "Calibri"
            r.font.size  = Pt(11) if ri == 0 else Pt(12)
            r.font.bold  = (ri == 0) or (ci == 0)
            r.font.color.rgb = ORANGE if ri == 0 else (WHITE if ci == 0 else TEXT)
            p.alignment = PP_ALIGN.LEFT

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl, BG)

# Right dark panel (42% of width)
PX = Inches(7.80)
PW = W - PX
rect(sl, PX, 0, PW, H, BG_DARK)
rect(sl, PX, 0, PW, Inches(0.10), ORANGE)
rect(sl, PX, 0, Inches(0.025), H, ORANGE)

# 5 stats stacked in right panel
stats5 = [("8","Containers"),("3","Backend Services"),
          ("20","User Stories"),("1","docker compose up"),("5","Days")]
cell_h = (H - Inches(0.10)) / 5
for i, (num, lbl) in enumerate(stats5):
    cy = Inches(0.10) + i * cell_h
    if i > 0:
        rect(sl, PX + Inches(0.35), cy, PW - Inches(0.7), Inches(0.018), BORDER)
    txt(sl, PX, cy + Inches(0.14), PW, Inches(0.62),
        num, sz=Pt(44), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txt(sl, PX, cy + Inches(0.80), PW, Inches(0.30),
        lbl, sz=Pt(11), color=MUTED, align=PP_ALIGN.CENTER)

# Left: orange left stripe + top bar
rect(sl, 0, 0, Inches(0.17), H, ORANGE)
rect(sl, 0, 0, PX, Inches(0.10), ORANGE)

# Decorative large faint rectangle (depth layer)
rect(sl, Inches(0.17), Inches(0.10), PX - Inches(0.17), Inches(1.35), BG_CARD)

# Event label
txt(sl, Inches(0.48), Inches(0.22), PX - Inches(0.6), Inches(0.32),
    "LAP 2025/2026  ·  HACKATHON — MARCH 2026",
    sz=Pt(10), color=MUTED)

# Thin line below event label
rect(sl, Inches(0.48), Inches(0.56), Inches(5.5), Inches(0.025), ORANGE)

# Giant title
txt(sl, Inches(0.48), Inches(1.62), Inches(7.15), Inches(1.52),
    "🔴  Mars Operations",
    sz=Pt(54), bold=True, color=ORANGE)

# Subtitle
txt(sl, Inches(0.48), Inches(3.24), Inches(7.15), Inches(0.60),
    "Distributed IoT Automation Platform",
    sz=Pt(22), color=TEXT)

# Orange divider
rect(sl, Inches(0.48), Inches(3.94), Inches(5.6), Inches(0.045), ORANGE)

# Description tag
txt(sl, Inches(0.48), Inches(4.08), Inches(7.15), Inches(0.38),
    "Event-driven  ·  Real-time SSE  ·  IF-THEN rules  ·  Fully containerised",
    sz=Pt(12.5), color=MUTED, italic=True)

# Pill: Student ID
rect(sl, Inches(0.48), Inches(4.68), Inches(2.52), Inches(0.46), BG_CARD, ORANGE, Pt(1.2))
rect(sl, Inches(0.48), Inches(4.68), Inches(0.09), Inches(0.46), ORANGE)
txt(sl, Inches(0.66), Inches(4.74), Inches(2.32), Inches(0.34),
    "Student ID: 2082600", sz=Pt(13), bold=True, color=ORANGE)

# Pill: Project
rect(sl, Inches(3.18), Inches(4.68), Inches(2.05), Inches(0.46), BG_CARD, ORANGE, Pt(1.2))
rect(sl, Inches(3.18), Inches(4.68), Inches(0.09), Inches(0.46), ORANGE)
txt(sl, Inches(3.36), Inches(4.74), Inches(1.87), Inches(0.34),
    "Project: MarsTok", sz=Pt(13), bold=True, color=ORANGE)

# Bottom orange footer
rect(sl, 0, H - Inches(0.30), PX, Inches(0.30), ORANGE)
txt(sl, Inches(0.48), H - Inches(0.28), PX - Inches(0.55), Inches(0.24),
    "Sapienza University of Rome  ·  A.Y. 2025/2026", sz=Pt(10), color=WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — PROBLEM & SOLUTION
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl, BG)
header(sl, "🚨  Problem & Solution")
footer(sl, 2)

rect(sl, Inches(0.5), CTOP, Inches(12.33), Inches(0.66), BG_CARD, BORDER)
rect(sl, Inches(0.5), CTOP, Inches(0.09), Inches(0.66), ORANGE)
txt(sl, Inches(0.70), CTOP + Inches(0.08), Inches(12.1), Inches(0.52),
    '"Mars 2036. Your automation stack is partially destroyed. Devices speak incompatible '
    'dialects. Rebuild it — or face thermodynamic consequences."',
    sz=Pt(12.5), italic=True, color=MUTED)

QH   = Inches(0.66)
COL  = Inches(6.0)
GAP  = Inches(0.33)
RX2  = Inches(0.5) + COL + GAP

txt(sl, Inches(0.5), CTOP+QH+Inches(0.12), COL, Inches(0.34),
    "❌  The Challenge", sz=Pt(14.5), bold=True, color=RED)
txt(sl, RX2,          CTOP+QH+Inches(0.12), COL, Inches(0.34),
    "✅  Our Solution",  sz=Pt(14.5), bold=True, color=GREEN)

challenges = [
    "15 devices,  8 raw JSON schemas,  two transport protocols",
    "REST polling + persistent SSE streams — no unified format",
    "Operators face a blank dashboard on page load",
    "Automation must fire without human intervention",
    "Configuration must survive service restarts",
]
solutions = [
    "Unified  InternalEvent  normalisation layer",
    "Event-driven pipeline via  RabbitMQ  fanout exchange",
    "IF-THEN rule engine — auto-triggers actuators",
    "React dashboard with  live SSE push",
    "Full Docker Compose IaC — one command start",
]
CH = Inches(0.72)
Y0 = CTOP + QH + Inches(0.54)
for i, (ch, sol) in enumerate(zip(challenges, solutions)):
    y = Y0 + i * (CH + Inches(0.07))
    card(sl, Inches(0.5), y, COL, CH, RED)
    txt(sl, Inches(0.70), y+Inches(0.17), Inches(5.75), Inches(0.50), ch, sz=Pt(12.5), color=TEXT)
    card(sl, RX2, y, COL, CH, GREEN)
    txt(sl, RX2+Inches(0.20), y+Inches(0.17), Inches(5.75), Inches(0.50), sol, sz=Pt(12.5), color=TEXT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — SYSTEM ARCHITECTURE
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl, BG)
header(sl, "🏗️  System Architecture")
footer(sl, 3)

LW3 = Inches(7.55)
LX3 = Inches(0.5)
BH3 = Inches(0.64)
AH3 = Inches(0.30)

arch = [
    ("🔴  mars-iot-simulator :8080",
     "8 REST sensors (poll 5 s)  ·  7 SSE telemetry topics", RED),
    ("🔵  collector  — Python 3.12 asyncio",
     "polls + streams + normalises (8 raw schemas) + publishes to RabbitMQ", BLUE),
    ("🟣  RabbitMQ 3.13  — fanout exchange  mars.events",
     "durable queues  ·  persistent delivery mode", PURPLE),
    ("🟠  rules-engine  — Python 3.12 asyncio",
     "cache state in Redis → evaluate rules → POST actuator → INSERT alert → PUBLISH", ORANGE),
]
arrows3 = [
    "↓  HTTP GET every 5 s   +   SSE EventSource (7 topics)",
    "↓  AMQP  PERSISTENT  publish  to  mars.events",
    "↓  exclusive anonymous queue  (auto-delete on consumer restart)",
    "↓  REST  ·  ↓  Redis state  ·  ↓  PostgreSQL  ·  ↓  Redis pub/sub",
]
y3 = CTOP
for i, (title, sub, acc) in enumerate(arch):
    card(sl, LX3, y3, LW3, BH3, acc)
    txt(sl, LX3+Inches(0.18), y3+Inches(0.06), LW3-Inches(0.25), Inches(0.28),
        title, sz=Pt(12.5), bold=True, color=WHITE)
    txt(sl, LX3+Inches(0.18), y3+Inches(0.36), LW3-Inches(0.25), Inches(0.24),
        sub, sz=Pt(10.5), color=MUTED)
    y3 += BH3
    txt(sl, LX3, y3, LW3, AH3, arrows3[i], sz=Pt(10), color=ORANGE, align=PP_ALIGN.CENTER)
    y3 += AH3

# api + frontend side by side
HW = (LW3 - Inches(0.12)) / 2
card(sl, LX3,             y3, HW, BH3, CYAN)
txt(sl, LX3+Inches(0.18), y3+Inches(0.08), HW-Inches(0.25), Inches(0.50),
    "🩵  api  — FastAPI :8000\nREST + SSE gateway  ·  pydantic v2", sz=Pt(11.5), color=WHITE)
card(sl, LX3+HW+Inches(0.12), y3, HW, BH3, GREEN)
txt(sl, LX3+HW+Inches(0.30), y3+Inches(0.08), HW-Inches(0.25), Inches(0.50),
    "🟢  frontend  — React 18 :3000\n7-page SPA  ·  Vite  ·  TailwindCSS", sz=Pt(11.5), color=WHITE)

# Right: infra sidebar
RX3 = Inches(8.42)
RW3 = Inches(4.41)
txt(sl, RX3, CTOP, RW3, Inches(0.32), "Infrastructure", sz=Pt(14), bold=True, color=ORNGL)

infra3 = [
    ("📨  RabbitMQ 3.13",  "fanout  mars.events\ndurable  ·  persistent delivery", PURPLE),
    ("⚡  Redis 7",         "state:{id}  TTL=3600 s\npub/sub:  mars.events  +  mars.alerts", CYAN),
    ("🐘  PostgreSQL 16",  "tables:  rules  +  alerts\nJSONB  ·  Alembic auto-migration", BLUE),
]
iy3 = CTOP + Inches(0.40)
for name, detail, acc in infra3:
    card(sl, RX3, iy3, RW3, Inches(1.82), acc)
    txt(sl, RX3+Inches(0.18), iy3+Inches(0.10), RW3-Inches(0.25), Inches(0.32),
        name, sz=Pt(12.5), bold=True, color=WHITE)
    txt(sl, RX3+Inches(0.18), iy3+Inches(0.50), RW3-Inches(0.25), Inches(1.2),
        detail, sz=Pt(11), color=MUTED)
    iy3 += Inches(1.82) + Inches(0.12)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — DATA PIPELINE & UNIFIED SCHEMA
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl, BG)
header(sl, "⚡  Data Pipeline & Unified Schema")
footer(sl, 4)

txt(sl, Inches(0.5), CTOP, Inches(5.9), Inches(0.32),
    "7-Step Pipeline", sz=Pt(13.5), bold=True, color=ORNGL)

steps4 = [
    ("Ingest",    "Poll 8 REST sensors every 5 s · maintain 7 persistent SSE connections"),
    ("Normalise", "8 raw schemas → 1 InternalEvent via dispatcher pattern"),
    ("Publish",   "PERSISTENT delivery to RabbitMQ fanout  mars.events"),
    ("Process",   "Write state:{id} to Redis TTL=1 h · evaluate all enabled rules"),
    ("Alert",     "POST actuator to simulator + INSERT PostgreSQL + PUBLISH Redis"),
    ("Stream",    "API subscribes Redis pub/sub → relays named SSE events to browser"),
    ("Render",    "useSSE: pre-fetch GET /api/state/ on load then merge live stream"),
]
SY4 = CTOP + Inches(0.40)
for i, (label, desc) in enumerate(steps4):
    y = SY4 + i * Inches(0.76)
    rect(sl, Inches(0.5), y, Inches(0.34), Inches(0.34), ORANGE)
    txt(sl, Inches(0.5), y+Inches(0.03), Inches(0.34), Inches(0.29),
        str(i+1), sz=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, Inches(0.95), y, Inches(1.30), Inches(0.34), label, sz=Pt(12), bold=True, color=WHITE)
    txt(sl, Inches(2.30), y, Inches(4.10), Inches(0.34), desc, sz=Pt(11.5), color=MUTED)
    if i < len(steps4)-1:
        rect(sl, Inches(0.5), y+Inches(0.38), Inches(5.88), Inches(0.01), BORDER)

RX4 = Inches(7.0)
RW4 = Inches(5.83)
txt(sl, RX4, CTOP, RW4, Inches(0.32),
    "InternalEvent Schema", sz=Pt(13.5), bold=True, color=ORNGL)

card(sl, RX4, CTOP+Inches(0.38), RW4, Inches(4.38))
schema4 = ('{\n'
           '  "event_id":    "550e8400-e29b-41d4-...",\n'
           '  "timestamp":   "2026-03-06T12:00:00Z",\n'
           '  "source_id":   "greenhouse_temperature",\n'
           '  "source_type": "rest_sensor",\n'
           '  "category":    "environment",\n'
           '  "metrics": [\n'
           '    { "name": "value", "value": 22.5, "unit": "degC" }\n'
           '  ],\n'
           '  "status":     "ok | warning | null",\n'
           '  "raw_schema": "rest.scalar.v1"\n'
           '}')
txt(sl, RX4+Inches(0.18), CTOP+Inches(0.50), RW4-Inches(0.25), Inches(4.18),
    schema4, sz=Pt(10.5), color=CODE, font="Courier New")

txt(sl, RX4, CTOP+Inches(4.88), RW4, Inches(0.26),
    "8 schemas handled:", sz=Pt(11), bold=True, color=WHITE)
txt(sl, RX4, CTOP+Inches(5.16), RW4, Inches(0.48),
    "rest.scalar.v1 · rest.chemistry.v1 · rest.level.v1 · rest.particulate.v1\n"
    "topic.power.v1 · topic.environment.v1 · topic.thermal_loop.v1 · topic.airlock.v1",
    sz=Pt(10.5), color=MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — AUTOMATION ENGINE
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl, BG)
header(sl, "⚙️  Automation Engine — IF-THEN Rules")
footer(sl, 5)

txt(sl, Inches(0.5), CTOP, Inches(6.1), Inches(0.32),
    "Rule Model (persisted in PostgreSQL)", sz=Pt(13.5), bold=True, color=ORNGL)

card(sl, Inches(0.5), CTOP+Inches(0.38), Inches(6.1), Inches(2.92))
rule_json = ('{\n'
             '  "name": "High Temp → Cooling Fan ON",\n'
             '  "enabled": true,\n'
             '  "condition": {\n'
             '    "source_id": "greenhouse_temperature",\n'
             '    "metric":    "value",\n'
             '    "operator":  ">",\n'
             '    "threshold": 35.0\n'
             '  },\n'
             '  "action": {\n'
             '    "actuator_name": "cooling_fan",\n'
             '    "state":         "ON"\n'
             '  }\n'
             '}')
txt(sl, Inches(0.68), CTOP+Inches(0.48), Inches(5.88), Inches(2.74),
    rule_json, sz=Pt(11.5), color=CODE, font="Courier New")

card(sl, Inches(0.5), CTOP+Inches(3.42), Inches(6.1), Inches(0.62), ORANGE)
txt(sl, Inches(0.68), CTOP+Inches(3.52), Inches(5.88), Inches(0.48),
    "🔁  Evaluated on every incoming event\n"
    "     SELECT * FROM rules WHERE enabled = true", sz=Pt(12), color=WHITE)

card(sl, Inches(0.5), CTOP+Inches(4.16), Inches(6.1), Inches(0.56), CYAN)
txt(sl, Inches(0.68), CTOP+Inches(4.26), Inches(5.88), Inches(0.38),
    'POST  simulator:8080/actuators/{name}  →  {"state": "ON"}',
    sz=Pt(11), color=WHITE, font="Courier New")

RX5 = Inches(7.1)
RW5 = Inches(5.73)
txt(sl, RX5, CTOP, RW5, Inches(0.32),
    "Supported Operators", sz=Pt(13.5), bold=True, color=ORNGL)

for i, op in enumerate([">  (gt)", "<  (lt)", ">=  (gte)", "<=  (lte)", "=  (eq)", "!=  (neq)"]):
    ci, ri = i % 3, i // 3
    bx = RX5 + ci * Inches(1.93)
    by = CTOP + Inches(0.40) + ri * Inches(0.57)
    rect(sl, bx, by, Inches(1.76), Inches(0.46), BG_CARD, BORDER)
    rect(sl, bx, by, Inches(0.07), Inches(0.46), ORANGE)
    txt(sl, bx+Inches(0.16), by+Inches(0.08), Inches(1.56), Inches(0.32),
        op, sz=Pt(13), bold=True, color=ORANGE)

txt(sl, RX5, CTOP+Inches(1.68), RW5, Inches(0.32),
    "Rule Lifecycle — Full CRUD", sz=Pt(13.5), bold=True, color=ORNGL)
add_table(sl, RX5, CTOP+Inches(2.08), RW5,
          [["Action","Endpoint","US"],
           ["Create",    "POST   /api/rules",             "13"],
           ["List/Get",  "GET    /api/rules",             "14"],
           ["Edit",      "PUT    /api/rules/{id}",        "15"],
           ["Toggle",    "PATCH  /api/rules/{id}/toggle", "16"],
           ["Delete",    "DELETE /api/rules/{id}",        "17"]],
          [Inches(1.15), Inches(3.50), Inches(0.65)], rh=Inches(0.38))

card(sl, RX5, CTOP+Inches(4.40), RW5, Inches(0.68), GREEN)
txt(sl, RX5+Inches(0.18), CTOP+Inches(4.52), RW5-Inches(0.25), Inches(0.50),
    "✅  Persisted in PostgreSQL — survive  docker compose down\n"
    "     alembic upgrade head  runs automatically on boot",
    sz=Pt(11.5), color=WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — FRONTEND DASHBOARD
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl, BG)
header(sl, "🖥️  Frontend Dashboard — 7 Pages")
footer(sl, 6)

add_table(sl, Inches(0.5), CTOP, Inches(12.33),
          [["Page",            "Key Features",                                                "US"],
           ["Dashboard",       "All sensor cards · live SSE · ok/warning status badges",     "1–4"],
           ["Power",           "6 live Recharts line charts (solar array, bus, consumption)", "5"],
           ["Environment",     "Radiation + life support cards · REST sensor widgets",        "6,9,10"],
           ["Airlock&Thermal", "State badge IDLE/PRESSURIZING/DEPRESSURIZING · thermal charts","7,8"],
           ["Actuators",       "ON/OFF toggle cards for all 4 actuators · 5 s auto-refresh", "11,12"],
           ["Rules",           "RuleForm dialog + RuleTable · full CRUD + toggle",            "13–18"],
           ["Alerts",          "Timeline · rule/source dropdowns · SSE live prepend",         "19,20"]],
          [Inches(1.62), Inches(9.32), Inches(0.96)], rh=Inches(0.44))

txt(sl, Inches(0.5), CTOP+Inches(4.02), Inches(12.33), Inches(0.30),
    "Real-time Data Flow", sz=Pt(13.5), bold=True, color=ORNGL)

flow6 = [
    ("Page load",             WHITE),
    ("GET /api/state/\npre-fetch all", CYAN),
    ("EventSource\n/api/stream SSE", ORANGE),
    ("sensor_update\nstate map", GREEN),
    ("alert event\nprepend list", AMBER),
    ("3 s reconnect\non drop", BLUE),
]
FBW6, FBH6 = Inches(1.90), Inches(0.82)
FY6 = CTOP + Inches(4.40)
for i, (label, col) in enumerate(flow6):
    fx = Inches(0.5) + i * (FBW6 + Inches(0.19))
    rect(sl, fx, FY6, FBW6, FBH6, BG_CARD, BORDER)
    rect(sl, fx, FY6, Inches(0.07), FBH6, col)
    txt(sl, fx+Inches(0.14), FY6+Inches(0.10), FBW6-Inches(0.18), FBH6-Inches(0.12),
        label, sz=Pt(10.5), color=col, align=PP_ALIGN.CENTER)
    if i < len(flow6)-1:
        txt(sl, fx+FBW6+Inches(0.01), FY6+Inches(0.27), Inches(0.17), Inches(0.30),
            "→", sz=Pt(16), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

txt(sl, Inches(0.5), CTOP+Inches(5.38), Inches(12.33), Inches(0.26),
    "React 18 · TypeScript · Vite · TailwindCSS · Recharts · lucide-react · react-router-dom v6",
    sz=Pt(10.5), color=MUTED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — TECHNOLOGY STACK
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl, BG)
header(sl, "🔧  Technology Stack")
footer(sl, 7)

add_table(sl, Inches(0.5), CTOP, Inches(12.33),
          [["Layer",          "Technology"],
           ["Ingestion",      "Python 3.12 · asyncio · httpx · aio-pika · pydantic v2"],
           ["Message Broker", "RabbitMQ 3.13 · fanout exchange mars.events · durable · persistent"],
           ["Processing",     "Python 3.12 · SQLAlchemy 2 async · asyncpg · alembic · httpx"],
           ["State Cache",    "Redis 7 · state:{id} TTL=3600 s · pub/sub: mars.events + mars.alerts"],
           ["Database",       "PostgreSQL 16 · rules + alerts tables · JSONB columns"],
           ["API Gateway",    "FastAPI · Uvicorn · sse-starlette · CORS middleware · pydantic v2"],
           ["Frontend",       "React 18 · TypeScript · Vite · TailwindCSS · Recharts · lucide-react"],
           ["IaC / Serving",  "Docker Compose · nginx · 8 containers · named volumes · healthchecks"]],
          [Inches(2.0), Inches(10.33)], rh=Inches(0.44))

stats7 = [("8","Containers"),("3","Backend\nServices"),("20","User Stories"),
          ("1","docker\ncompose up"),("5","Days")]
SW7 = Inches(12.33) / len(stats7)
for i, (num, lbl) in enumerate(stats7):
    sx = Inches(0.5) + i * SW7
    rect(sl, sx+Inches(0.08), CTOP+Inches(4.22), SW7-Inches(0.16), Inches(1.22), BG_CARD, BORDER)
    rect(sl, sx+Inches(0.08), CTOP+Inches(4.22), SW7-Inches(0.16), Inches(0.06), ORANGE)
    txt(sl, sx+Inches(0.08), CTOP+Inches(4.32), SW7-Inches(0.16), Inches(0.52),
        num, sz=Pt(32), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txt(sl, sx+Inches(0.08), CTOP+Inches(4.90), SW7-Inches(0.16), Inches(0.40),
        lbl, sz=Pt(10), color=MUTED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — INFRASTRUCTURE AS CODE
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl, BG)
header(sl, "🐳  Infrastructure as Code")
footer(sl, 8)

LW8 = Inches(6.25)
txt(sl, Inches(0.5), CTOP, LW8, Inches(0.32),
    "One Command Start", sz=Pt(13.5), bold=True, color=ORNGL)

card(sl, Inches(0.5), CTOP+Inches(0.38), LW8, Inches(1.42))
bash8 = ("# One-time: load the OCI simulator image\n"
         "./source/load-image.sh\n\n"
         "# Start the entire platform\n"
         "cd source && docker compose up")
txt(sl, Inches(0.68), CTOP+Inches(0.48), LW8-Inches(0.25), Inches(1.26),
    bash8, sz=Pt(12.5), color=CODE, font="Courier New")

txt(sl, Inches(0.5), CTOP+Inches(1.94), LW8, Inches(0.32),
    "Startup Dependency Chain", sz=Pt(13.5), bold=True, color=ORNGL)

chain8 = [
    ("🟣 rabbitmq  ·  ⚡ redis  ·  🐘 postgres  ·  🔴 simulator", PURPLE),
    ("🔵 collector   ·   🟠 rules-engine",                          BLUE),
    ("🩵 api  :8000",                                               CYAN),
    ("🟢 frontend  :3000",                                          GREEN),
]
cy8 = CTOP + Inches(2.34)
for j, (lbl, acc) in enumerate(chain8):
    card(sl, Inches(0.5), cy8, LW8, Inches(0.58), acc)
    txt(sl, Inches(0.68), cy8+Inches(0.13), LW8-Inches(0.25), Inches(0.36),
        lbl, sz=Pt(12), color=WHITE, align=PP_ALIGN.CENTER)
    cy8 += Inches(0.58)
    if j < len(chain8)-1:
        txt(sl, Inches(2.5), cy8, Inches(1.5), Inches(0.22),
            "↓  condition: service_healthy", sz=Pt(9), color=MUTED, align=PP_ALIGN.CENTER)
        cy8 += Inches(0.22)

RX8 = Inches(7.1)
RW8 = Inches(5.73)
txt(sl, RX8, CTOP, RW8, Inches(0.32),
    "Named Volumes & Persistence", sz=Pt(13.5), bold=True, color=ORNGL)

vols = [
    ("🗄️  pg_data",       "Rules + alerts persist across restarts\nUS 18 — rule persistence verified", ORANGE),
    ("📨  rabbitmq_data",  "Durable queue messages survive broker restart", ORANGE),
    ("⚡  redis_data",     "State cache repopulates from RabbitMQ stream\nwithin seconds on loss", ORANGE),
]
vy8 = CTOP+Inches(0.38)
for name, desc, acc in vols:
    card(sl, RX8, vy8, RW8, Inches(1.30), acc)
    txt(sl, RX8+Inches(0.18), vy8+Inches(0.10), RW8-Inches(0.25), Inches(0.32),
        name, sz=Pt(13), bold=True, color=WHITE)
    txt(sl, RX8+Inches(0.18), vy8+Inches(0.50), RW8-Inches(0.25), Inches(0.72),
        desc, sz=Pt(11.5), color=MUTED)
    vy8 += Inches(1.30)+Inches(0.14)

card(sl, RX8, vy8+Inches(0.08), RW8, Inches(0.60), ORANGE)
txt(sl, RX8+Inches(0.18), vy8+Inches(0.18), RW8-Inches(0.25), Inches(0.44),
    "No manual steps after  docker compose up\n"
    "Alembic upgrade head runs automatically on boot",
    sz=Pt(12), color=WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — LESSONS LEARNED & DESIGN DECISIONS
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl, BG)
header(sl, "💡  Lessons Learned & Design Decisions")
footer(sl, 9)

COL9 = Inches(6.0)
RX9  = Inches(0.5)+COL9+Inches(0.33)

txt(sl, Inches(0.5), CTOP, COL9, Inches(0.32),
    "⚠️  Engineering Challenges", sz=Pt(13.5), bold=True, color=AMBER)
txt(sl, RX9, CTOP, COL9, Inches(0.32),
    "✅  Key Design Decisions", sz=Pt(13.5), bold=True, color=GREEN)

chal9 = [
    ("Schema diversity",
     "8 raw device formats required a dispatcher pattern: each schema maps to "
     "its own handler in the collector normaliser."),
    ("SSE double reconnection",
     "Both collector (→ simulator) and frontend (→ API) need independent 3 s "
     "auto-retry loops to survive drops."),
    ("Blank dashboard on load",
     "Opening SSE first = empty cards for up to 5 s. Fix: pre-fetch "
     "GET /api/state/ before opening EventSource."),
    ("JSONB alert filtering",
     "Used PostgreSQL @> containment operator to filter source_id inside "
     "the stored JSONB event blob."),
]
decs9 = [
    ("Fanout exchange",
     "Collector publishes once → any new service binds its own anonymous queue. "
     "Zero collector changes to extend the system."),
    ("Redis pub/sub relay",
     "True zero-latency push from rules-engine to API SSE clients — "
     "no polling, no sleep loops needed."),
    ("Alembic on boot",
     "alembic upgrade head at rules-engine startup: DB schema always in sync, "
     "no manual migration step ever required."),
    ("Exclusive anonymous queues",
     "Auto-deleted on consumer restart: no stale messages accumulate. "
     "Broker stays clean automatically."),
]
CH9 = Inches(1.22)
y9  = CTOP+Inches(0.40)
for i, ((ct, cd), (dt, dd)) in enumerate(zip(chal9, decs9)):
    y = y9 + i*(CH9+Inches(0.10))
    card(sl, Inches(0.5), y, COL9, CH9, AMBER)
    txt(sl, Inches(0.68), y+Inches(0.08), COL9-Inches(0.25), Inches(0.32),
        ct, sz=Pt(12.5), bold=True, color=WHITE)
    txt(sl, Inches(0.68), y+Inches(0.46), COL9-Inches(0.25), Inches(0.70),
        cd, sz=Pt(11), color=MUTED)
    card(sl, RX9, y, COL9, CH9, GREEN)
    txt(sl, RX9+Inches(0.18), y+Inches(0.08), COL9-Inches(0.25), Inches(0.32),
        dt, sz=Pt(12.5), bold=True, color=WHITE)
    txt(sl, RX9+Inches(0.18), y+Inches(0.46), COL9-Inches(0.25), Inches(0.70),
        dd, sz=Pt(11), color=MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — LIVE DEMO & CLOSING
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl, BG)
header(sl, "🎬  Live Demo")
footer(sl, 10)

txt(sl, Inches(0.5), CTOP, Inches(7.5), Inches(0.32),
    "Demo Script", sz=Pt(13.5), bold=True, color=ORNGL)

demo = [
    ("docker compose up",  "— all 8 containers reach healthy state in dependency order"),
    ("Dashboard",          "— cards populate instantly (US 4 pre-fetch); values live every 5 s"),
    ("Power page",         "— 6 Recharts live line charts with rolling window"),
    ("Actuators",          "— toggle  cooling_fan  ON → OFF; optimistic + 5 s auto-confirm"),
    ("Rules",              "— create:  greenhouse_temperature > 28 → cooling_fan ON  (US 13)"),
    ("Alerts page",        "— alert appears live via SSE; filter by rule / source_id"),
    ("Persistence",        "— down && up → rules still there ✅  (US 18 verified)"),
]
sy10 = CTOP+Inches(0.40)
for i, (label, desc) in enumerate(demo):
    y = sy10 + i*Inches(0.72)
    rect(sl, Inches(0.5), y+Inches(0.04), Inches(0.32), Inches(0.32), ORANGE)
    txt(sl, Inches(0.5), y+Inches(0.04), Inches(0.32), Inches(0.28),
        str(i+1), sz=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, Inches(0.94), y, Inches(1.85), Inches(0.36), label, sz=Pt(12), bold=True, color=WHITE)
    txt(sl, Inches(2.85), y, Inches(5.10), Inches(0.36), desc, sz=Pt(11.5), color=MUTED)

RX10 = Inches(8.5)
RW10 = Inches(4.38)
txt(sl, RX10, CTOP, RW10, Inches(0.32),
    "Endpoints", sz=Pt(13.5), bold=True, color=ORNGL)

eps = [
    ("🖥️  http://localhost:3000",     "React dashboard",                    CYAN),
    ("📖  http://localhost:8000/docs", "FastAPI Swagger / ReDoc",            BLUE),
    ("📨  http://localhost:15672",     "RabbitMQ Management  (guest/guest)", PURPLE),
]
ey10 = CTOP+Inches(0.40)
for url, label, col in eps:
    card(sl, RX10, ey10, RW10, Inches(0.90), col)
    txt(sl, RX10+Inches(0.18), ey10+Inches(0.08), RW10-Inches(0.25), Inches(0.32),
        url, sz=Pt(11.5), bold=True, color=col, font="Courier New")
    txt(sl, RX10+Inches(0.18), ey10+Inches(0.50), RW10-Inches(0.25), Inches(0.30),
        label, sz=Pt(11), color=MUTED)
    ey10 += Inches(0.90)+Inches(0.12)

stats10 = [("20","User Stories"),("8","Containers"),("5","Days")]
SW10 = RW10 / 3
for i, (num, lbl) in enumerate(stats10):
    sx = RX10 + i*SW10
    rect(sl, sx+Inches(0.06), ey10, SW10-Inches(0.12), Inches(0.90), BG_CARD, BORDER)
    txt(sl, sx+Inches(0.06), ey10+Inches(0.06), SW10-Inches(0.12), Inches(0.46),
        num, sz=Pt(26), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txt(sl, sx+Inches(0.06), ey10+Inches(0.58), SW10-Inches(0.12), Inches(0.26),
        lbl, sz=Pt(9.5), color=MUTED, align=PP_ALIGN.CENTER)

rect(sl, RX10, ey10+Inches(1.02), RW10, Inches(0.76), ORANGE)
txt(sl, RX10, ey10+Inches(1.08), RW10, Inches(0.64),
    "1 command  ·  0 manual setup\nNo blank screens.",
    sz=Pt(14.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
OUT = "/Users/davide/Desktop/Mars/presentation_editable.pptx"
prs.save(OUT)
print(f"Saved → {OUT}  ({len(prs.slides)} slides)")
